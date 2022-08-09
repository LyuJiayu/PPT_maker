import collections 
import collections.abc

from pptx import Presentation
from pptx.util import Cm
import os

filePath = os.getcwd()
my_files = os.listdir(filePath)

prs = Presentation('test.pptx') #打开一个PPT文件
prs.slide_height = Cm(19.05)    #设置ppt的高度
prs.slide_width = Cm(33.867)  #设置ppt的宽度

def m_tf(slide,left,top,width,height,text): #插入文本框，并写入文字
    left = Cm(left)
    top = Cm(top)
    width = Cm(width)
    height = Cm(height)
    txBox = slide.shapes.add_textbox(left,top,width,height)
    tf = txBox.text_frame
    tf.text = text
    return 0

def get_range(path):
    with open('{}'.format(path),'rb') as temp_f:
        datafile = temp_f.readlines()
        s = str(datafile[19]).split() # sxm文件第19行为扫描尺寸[m]
        f = float(s[1])*1000000000 # 单位换算[m]->[nm]
        return '   {:.0f}nm'.format(f)
            
def get_text(path):
    with open('{}'.format(path),'rb') as temp_f:
        datafile = temp_f.readlines()
        s = str(datafile[15]).split('\\') # sxm文件第15行为文件名
        s = s[-2].split('.')
        s = s[0]
        return s

my_range = []
my_text = []
x = input('你要读取的sxm文件的前缀为?')
for file in my_files:
    if x in file:
        my_range.append(get_range(file))
        my_text.append(get_text(file))


for i in range(len(my_text)):
    if i%6 == 0: # 每6个sxm文件新建一张PPT
        blank_slide_layout = prs.slide_layouts[6] #选择母版6
        slide = prs.slides.add_slide(blank_slide_layout) #添加新一页幻灯片
    if i%6 == 0: m_tf(slide,4.5,8,8,1,my_text[i]+my_range[i])
    if i%6 == 1: m_tf(slide,13.5,8,8,1,my_text[i]+my_range[i])
    if i%6 == 2: m_tf(slide,23,8,8,1,my_text[i]+my_range[i])
    if i%6 == 3: m_tf(slide,4.5,17,8,1,my_text[i]+my_range[i])
    if i%6 == 4: m_tf(slide,13.5,17,8,1,my_text[i]+my_range[i])
    if i%6 == 5: m_tf(slide,23,17,8,1,my_text[i]+my_range[i])

prs.save('test.pptx')